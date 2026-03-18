"""
Document processor abstraction for llmLibrarian ingestion.
Each processor handles extraction for a specific file type.
"""
import hashlib
import io
import importlib.util
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Protocol

# ChunkTuple imported from ingest to avoid circular; use forward reference in type hints.
ChunkTuple = tuple[str, str, dict[str, Any]]


@dataclass(frozen=True)
class ExtractedText:
    text: str
    meta: dict[str, Any] | None = None


@dataclass(frozen=True)
class ExtractedPage:
    text: str
    page_num: int
    meta: dict[str, Any] | None = None


@dataclass(frozen=True)
class ImageRegion:
    text: str
    role: str
    x: float
    y: float
    w: float
    h: float
    needs_vision_enrichment: bool = False
    meta: dict[str, Any] | None = None


@dataclass(frozen=True)
class ExtractedImage:
    summary: str
    visible_text: str
    regions: tuple[ImageRegion, ...]
    meta: dict[str, Any] | None = None
    artifact: dict[str, Any] | None = None


class DocumentProcessor(Protocol):
    """Protocol for document processors. Each handles one file type."""

    def extract(self, data: bytes, source_path: str) -> str | ExtractedText | ExtractedImage | list[ExtractedPage] | None:
        """Extract text from file bytes. Returns:
        - str/ExtractedText for plain text content
        - ExtractedImage for image files with OCR/vision enrichment
        - list[ExtractedPage] for paged documents (PDF)
        - None if extraction failed or library unavailable
        """
        ...


class DocumentExtractionError(Exception):
    """Base error for document extraction failures."""


class PDFExtractionError(DocumentExtractionError):
    """Raised when PDF extraction fails."""


class DOCXExtractionError(DocumentExtractionError):
    """Raised when DOCX extraction fails."""


class XLSXExtractionError(DocumentExtractionError):
    """Raised when XLSX extraction fails."""


class PPTXExtractionError(DocumentExtractionError):
    """Raised when PPTX extraction fails."""


class ImageExtractionError(DocumentExtractionError):
    """Raised when image OCR extraction fails."""


class TextExtractionError(DocumentExtractionError):
    """Raised when text extraction fails."""

    @property
    def format_label(self) -> str:
        """Human-readable format name for metadata-only fallback."""
        ...

    @property
    def install_hint(self) -> str:
        """Package name hint when extractor is unavailable."""
        ...


_PDF_STRUCTURED_MAX_CHARS = 6000
_NUMERIC_VALUE_PATTERN = re.compile(r"^\$?\s*-?\d[\d,]*(?:\.\d+)?\.?$")
_LINE_TOKEN_PATTERN = re.compile(r"^(?:line\s*)?(\d{1,3}[a-z]?)$", re.IGNORECASE)
_MAX_HINT_LINE_NUMBER = 80
_OCR_LABEL = "OCR text (scan fallback):"
_PADDLE_OCR_ENGINE: Any | None = None
_PADDLE_OCR_INIT_ATTEMPTED = False
_PADDLE_OCR_USE_ANGLE: bool | None = None
_OCR_PREPROCESS_WARNED: set[str] = set()
_PROCESSOR_LOG_LEVELS = {"DEBUG": 10, "INFO": 20, "WARN": 30, "ERROR": 40}
_VALID_OCR_BACKENDS = frozenset({"auto", "vision", "paddleocr", "tesseract"})
_VISION_HELPER_BINARY: Path | None = None
_VISION_HELPER_READY: bool | None = None
_VISION_MODEL_READY: dict[str, bool] = {}
_IMAGE_EXTENSIONS = (".png", ".jpg", ".jpeg", ".heic", ".heif", ".tif", ".tiff")


@dataclass(frozen=True)
class _OCRResult:
    text: str
    backend: str
    observations: tuple[dict[str, Any], ...] = ()
    raw_payload: dict[str, Any] | None = None
    quality_reasons: tuple[str, ...] = ()


def _processor_min_log_level() -> str:
    raw = (os.environ.get("LLMLIBRARIAN_PROCESSOR_LOG_LEVEL") or "WARN").strip().upper()
    if raw == "WARNING":
        raw = "WARN"
    return raw if raw in _PROCESSOR_LOG_LEVELS else "WARN"


def _log_processor_event(level: str, message: str, **fields: Any) -> None:
    """Write structured processor events in the same JSON-lines style as ingest logs."""
    normalized_level = str(level or "INFO").upper()
    min_level = _processor_min_log_level()
    if _PROCESSOR_LOG_LEVELS.get(normalized_level, 20) < _PROCESSOR_LOG_LEVELS[min_level]:
        return
    payload = {"ts": datetime.now(timezone.utc).isoformat(), "level": normalized_level, "message": message}
    if fields:
        payload.update(fields)
    try:
        print(json.dumps(payload, ensure_ascii=False), file=sys.stderr)
    except Exception:
        pass


def _log_warn_once(key: str, message: str, **fields: Any) -> None:
    """Emit a warning once per process for noisy repeated failures."""
    if key in _OCR_PREPROCESS_WARNED:
        return
    _OCR_PREPROCESS_WARNED.add(key)
    _log_processor_event("WARN", message, **fields)


def _pdf_tables_enabled() -> bool:
    # Default OFF: table extraction via pdfplumber is expensive and noisy on malformed/scanned PDFs.
    val = os.environ.get("LLMLIBRARIAN_PDF_TABLES", "0").strip().lower()
    return val not in ("0", "false", "no")


def _ocr_preprocess_enabled() -> bool:
    """Whether OCR preprocessing and enhanced OCR mode is enabled."""
    val = os.environ.get("LLMLIBRARIAN_OCR_PREPROCESS", "0").strip().lower()
    return val not in ("0", "false", "no")


def _configured_ocr_backend() -> str:
    raw = (os.environ.get("LLMLIBRARIAN_OCR_BACKEND") or "auto").strip().lower()
    return raw if raw in _VALID_OCR_BACKENDS else "auto"


def _configured_vision_model() -> str:
    return (os.environ.get("LLMLIBRARIAN_VISION_MODEL") or "").strip()


def _vision_helper_source_path() -> Path:
    return Path(__file__).with_name("vision_ocr.swift")


def _vision_prereqs_available() -> bool:
    return (
        sys.platform == "darwin"
        and bool(shutil.which("swiftc"))
        and _vision_helper_source_path().exists()
    )


def _vision_ocr_available() -> bool:
    if _VISION_HELPER_READY is not None:
        return _VISION_HELPER_READY
    return _vision_prereqs_available()


def _preferred_ocr_backends() -> list[str]:
    configured = _configured_ocr_backend()
    if configured != "auto":
        return [configured]
    ordered: list[str] = []
    if sys.platform == "darwin":
        ordered.append("vision")
    ordered.extend(["paddleocr", "tesseract"])
    return ordered


def _paddleocr_available() -> bool:
    """True when paddleocr is importable in the current environment."""
    try:
        return (
            importlib.util.find_spec("paddleocr") is not None
            and importlib.util.find_spec("paddle") is not None
        )
    except Exception:
        return False


def _tesseract_available() -> bool:
    """True when the tesseract binary is available on PATH."""
    return bool(shutil.which("tesseract"))


def _ocr_backend_available(name: str) -> bool:
    if name == "vision":
        return _vision_ocr_available()
    if name == "paddleocr":
        return _paddleocr_available()
    if name == "tesseract":
        return _tesseract_available()
    return False


def _available_ocr_backends() -> list[str]:
    """Return available OCR backends in deterministic priority order."""
    return [backend for backend in _preferred_ocr_backends() if _ocr_backend_available(backend)]


def ocr_backend_chain_for_capabilities() -> list[str]:
    """Backends that would be attempted in auto mode on this machine."""
    return _available_ocr_backends()


def _alnum_ratio(s: str) -> float:
    if not s:
        return 0.0
    alnum = sum(1 for c in s if c.isalnum() or c in " $.,")
    return alnum / len(s)


def _ocr_quality_assessment(text: str) -> tuple[bool, tuple[str, ...], dict[str, Any]]:
    """Conservative OCR quality gate tuned to reject symbol-heavy gibberish."""
    raw = (text or "").strip()
    if not raw:
        return False, ("empty_text",), {"length": 0}

    alnum_ratio = _alnum_ratio(raw)
    tokens = re.findall(r"[A-Za-z0-9$][A-Za-z0-9$'.,:-]*", raw)
    alpha_tokens = [tok for tok in tokens if re.search(r"[A-Za-z]", tok)]
    multi_char_alpha = [tok for tok in alpha_tokens if len(re.sub(r"[^A-Za-z]", "", tok)) >= 3]
    single_char_tokens = sum(1 for tok in tokens if len(re.sub(r"[^A-Za-z0-9]", "", tok)) == 1)
    symbol_chars = sum(1 for ch in raw if not (ch.isalnum() or ch.isspace() or ch in "$.,:;!?'-/()&"))
    punct_ratio = symbol_chars / len(raw) if raw else 1.0
    digit_chars = sum(1 for ch in raw if ch.isdigit())
    stats = {
        "length": len(raw),
        "token_count": len(tokens),
        "alpha_token_count": len(alpha_tokens),
        "multi_char_alpha_count": len(multi_char_alpha),
        "single_char_tokens": single_char_tokens,
        "alnum_ratio": round(alnum_ratio, 4),
        "punct_ratio": round(punct_ratio, 4),
        "digit_ratio": round((digit_chars / len(raw)) if raw else 0.0, 4),
    }

    reasons: list[str] = []
    if len(raw) >= 18 and alnum_ratio < 0.55:
        reasons.append("low_alnum_ratio")
    if len(raw) >= 25 and punct_ratio > 0.18:
        reasons.append("high_symbol_ratio")
    if len(tokens) >= 5 and len(multi_char_alpha) < 2 and digit_chars < 4:
        reasons.append("missing_multi_char_words")
    if len(tokens) >= 12 and single_char_tokens >= max(6, int(len(tokens) * 0.4)):
        reasons.append("too_many_single_char_tokens")
    if len(raw) >= 40 and len(alpha_tokens) >= 4 and len(multi_char_alpha) == 0:
        reasons.append("alpha_tokens_too_short")
    if len(alpha_tokens) >= 12 and len(multi_char_alpha) < max(3, int(len(alpha_tokens) * 0.35)):
        reasons.append("low_meaningful_word_ratio")
    if len(raw) < 12 and alnum_ratio >= 0.7:
        reasons = []

    return (not reasons, tuple(reasons), stats)


def _image_summary_placeholder(visible_text: str) -> str:
    cleaned = (visible_text or "").strip()
    if cleaned:
        return "Text-forward image with deferred visual summary."
    return "Photo image with deferred visual summary; no reliable OCR text."


def _image_summary_disabled_placeholder(visible_text: str) -> str:
    cleaned = (visible_text or "").strip()
    if cleaned:
        return "Text-forward image indexed with OCR only; multimodal vision disabled."
    return "Photo image indexed with OCR only; multimodal vision disabled."


def _image_ocr_signal_assessment(ocr_result: _OCRResult | None) -> dict[str, Any]:
    visible_text = (ocr_result.text if ocr_result else "").strip()
    observations = list(ocr_result.observations) if ocr_result else []
    backend = str(ocr_result.backend or "") if ocr_result else ""
    quality_ok, reasons, stats = _ocr_quality_assessment(visible_text)
    multi_char_alpha = int(stats.get("multi_char_alpha_count") or 0)
    token_count = int(stats.get("token_count") or 0)
    alpha_token_count = int(stats.get("alpha_token_count") or 0)
    single_char_tokens = int(stats.get("single_char_tokens") or 0)
    alnum_ratio = float(stats.get("alnum_ratio") or 0.0)
    text_len = len(visible_text)
    observation_count = len(observations)
    coverage = 0.0
    if observations:
        coverage = sum(
            max(0.0, min(1.0, float(o.get("w") or 0.0))) * max(0.0, min(1.0, float(o.get("h") or 0.0)))
            for o in observations
            if isinstance(o, dict)
        )
        coverage = min(1.0, coverage)
    single_char_ratio = (single_char_tokens / token_count) if token_count else 1.0
    meaningful_word_ratio = (multi_char_alpha / alpha_token_count) if alpha_token_count else 0.0

    score = 0.0
    if quality_ok:
        score += 0.3
    score += min(0.3, multi_char_alpha * 0.06)
    score += min(0.18, observation_count * 0.04)
    score += min(0.16, text_len / 300.0)
    score += min(0.16, coverage * 1.6)
    score = round(min(1.0, score), 4)

    text_structured = bool(
        quality_ok
        and backend == "vision"
        and (
            observation_count >= 3
            or (observation_count >= 2 and coverage >= 0.018 and multi_char_alpha >= 4 and text_len >= 30)
            or (observation_count >= 1 and coverage >= 0.06 and multi_char_alpha >= 4 and text_len >= 20)
        )
    )
    ocr_strong = bool(
        quality_ok
        and observation_count == 0
        and backend in {"tesseract", "paddleocr"}
        and multi_char_alpha >= 8
        and text_len >= 80
        and single_char_ratio <= 0.15
        and meaningful_word_ratio >= 0.7
        and alnum_ratio >= 0.85
    )
    eager_summary = bool(text_structured or ocr_strong)
    keep_visible_text = bool(
        visible_text
        and quality_ok
        and (
            score >= 0.4
            or multi_char_alpha >= 2
            or (observation_count >= 2 and text_len >= 18)
            or text_len < 16
        )
    )
    return {
        "ocr_signal_score": score,
        "quality_ok": quality_ok,
        "quality_reasons": reasons,
        "quality_stats": stats,
        "backend": backend,
        "observation_count": observation_count,
        "coverage": round(coverage, 4),
        "single_char_ratio": round(single_char_ratio, 4),
        "meaningful_word_ratio": round(meaningful_word_ratio, 4),
        "text_structured": text_structured,
        "ocr_strong": ocr_strong,
        "eager_summary": eager_summary,
        "keep_visible_text": keep_visible_text,
    }


def ensure_vision_model_ready() -> str:
    """
    Ensure the configured Ollama vision model exists and advertises image capability.
    Raises ImageExtractionError when image enrichment is required but unavailable.
    """
    model = _configured_vision_model()
    if not model:
        raise ImageExtractionError(
            "LLMLIBRARIAN_VISION_MODEL is required for standalone image enrichment."
        )
    if _VISION_MODEL_READY.get(model):
        return model
    try:
        import ollama

        info = ollama.show(model)
        capabilities = [str(v).lower() for v in (getattr(info, "capabilities", None) or [])]
    except Exception as e:
        raise ImageExtractionError(
            f"LLMLIBRARIAN_VISION_MODEL={model!r} is unavailable: {e}"
        ) from e
    if "vision" not in capabilities:
        raise ImageExtractionError(
            f"LLMLIBRARIAN_VISION_MODEL={model!r} is not vision-capable."
        )
    _VISION_MODEL_READY[model] = True
    return model


def _preprocess_image_for_ocr(image_bytes: bytes) -> bytes:
    """Apply OpenCV-based cleanup to scanned pages before OCR."""
    steps: list[str] = []
    try:
        try:
            import cv2
            import numpy as np
        except Exception as e:
            _log_warn_once(
                "ocr_preprocess_import",
                "OCR preprocessing unavailable; missing optional dependencies",
                error=str(e),
                extractor="opencv",
            )
            return image_bytes

        arr = np.frombuffer(image_bytes, dtype=np.uint8)
        image = cv2.imdecode(arr, cv2.IMREAD_COLOR)
        if image is None:
            _log_warn_once(
                "ocr_preprocess_decode",
                "OCR preprocessing decode failed; using original page image",
                extractor="opencv",
            )
            return image_bytes

        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        steps.append("grayscale")
        denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
        steps.append("denoise")
        blurred = cv2.GaussianBlur(denoised, (3, 3), 0)
        steps.append("gaussian_blur")
        _threshold, binarized = cv2.threshold(
            blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
        )
        steps.append("otsu_binarize")

        deskewed = binarized
        coords = np.column_stack(np.where(binarized < 128))
        if len(coords) > 5:
            rect = cv2.minAreaRect(coords.astype("float32"))
            raw_angle = float(rect[-1])
            angle = -(90.0 + raw_angle) if raw_angle < -45.0 else -raw_angle
            if 0.5 < abs(angle) < 15.0:
                h, w = binarized.shape[:2]
                center = (w / 2.0, h / 2.0)
                matrix = cv2.getRotationMatrix2D(center, angle, 1.0)
                deskewed = cv2.warpAffine(
                    binarized,
                    matrix,
                    (w, h),
                    flags=cv2.INTER_CUBIC,
                    borderMode=cv2.BORDER_REPLICATE,
                )
                steps.append(f"deskew({angle:.1f}deg)")

        ok, encoded = cv2.imencode(".png", deskewed)
        if not ok:
            _log_warn_once(
                "ocr_preprocess_encode",
                "OCR preprocessing encode failed; using original page image",
                extractor="opencv",
            )
            return image_bytes

        if steps:
            _log_processor_event("DEBUG", "OCR preprocessing applied", steps=steps)
        return encoded.tobytes()
    except Exception as e:
        _log_warn_once(
            "ocr_preprocess_error",
            "OCR preprocessing failed; using original page image",
            error=str(e),
            extractor="opencv",
        )
        return image_bytes


def _get_paddle_ocr_engine() -> Any | None:
    """Lazy-init PaddleOCR engine once per process; returns None when unavailable."""
    global _PADDLE_OCR_ENGINE, _PADDLE_OCR_INIT_ATTEMPTED, _PADDLE_OCR_USE_ANGLE
    use_angle = _ocr_preprocess_enabled()
    if _PADDLE_OCR_INIT_ATTEMPTED and _PADDLE_OCR_USE_ANGLE == use_angle:
        return _PADDLE_OCR_ENGINE
    _PADDLE_OCR_INIT_ATTEMPTED = True
    _PADDLE_OCR_USE_ANGLE = use_angle
    if not _paddleocr_available():
        _PADDLE_OCR_ENGINE = None
        return None
    try:
        # Avoid startup network checks in local/offline workflows.
        os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "True")
        from paddleocr import PaddleOCR

        try:
            _PADDLE_OCR_ENGINE = PaddleOCR(use_angle_cls=use_angle, lang="en", show_log=False)
        except Exception as e:
            # PaddleOCR versions differ on constructor kwargs (e.g., no show_log).
            if "show_log" not in str(e):
                raise
            _PADDLE_OCR_ENGINE = PaddleOCR(use_angle_cls=use_angle, lang="en")
    except Exception as e:
        _PADDLE_OCR_ENGINE = None
        _log_processor_event(
            "WARN",
            "PaddleOCR initialization failed",
            error=str(e),
            extractor="paddleocr",
        )
    return _PADDLE_OCR_ENGINE


def _extract_paddleocr_text(result: Any) -> str:
    """Extract OCR text fragments from PaddleOCR result payload."""
    parts: list[str] = []

    def _walk(node: Any) -> None:
        if not isinstance(node, (list, tuple)):
            return
        if len(node) >= 2 and isinstance(node[1], (list, tuple)) and node[1]:
            text_candidate = node[1][0]
            if isinstance(text_candidate, str):
                cleaned = text_candidate.strip()
                if cleaned:
                    parts.append(cleaned)
        for child in node:
            _walk(child)

    _walk(result)
    seen: set[str] = set()
    deduped: list[str] = []
    for part in parts:
        if part in seen:
            continue
        seen.add(part)
        deduped.append(part)
    return " ".join(deduped).strip()


def _vision_helper_binary_path() -> Path:
    source = _vision_helper_source_path()
    try:
        digest = hashlib.sha256(source.read_bytes()).hexdigest()[:12]
    except OSError:
        digest = "missing"
    cache_dir = Path(tempfile.gettempdir()) / "llmlibrarian-vision"
    return cache_dir / f"vision-ocr-{digest}"


def _ensure_vision_helper_binary() -> str | None:
    global _VISION_HELPER_BINARY, _VISION_HELPER_READY
    if _VISION_HELPER_READY and _VISION_HELPER_BINARY is not None and _VISION_HELPER_BINARY.exists():
        return str(_VISION_HELPER_BINARY)
    if not _vision_prereqs_available():
        _VISION_HELPER_READY = False
        _VISION_HELPER_BINARY = None
        return None

    source = _vision_helper_source_path()
    binary = _vision_helper_binary_path()
    binary.parent.mkdir(parents=True, exist_ok=True)
    if binary.exists():
        _VISION_HELPER_BINARY = binary
        _VISION_HELPER_READY = True
        _log_processor_event("DEBUG", "Vision OCR helper available", extractor="vision", binary=str(binary))
        return str(binary)

    tmp_binary = binary.with_suffix(".tmp")
    cmd = ["swiftc", "-O", str(source), "-o", str(tmp_binary)]
    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if proc.returncode != 0:
            _VISION_HELPER_READY = False
            _VISION_HELPER_BINARY = None
            _log_processor_event(
                "WARN",
                "Vision OCR helper compile failed",
                extractor="vision",
                error=(proc.stderr or proc.stdout or "").strip()[:400],
            )
            return None
        os.replace(tmp_binary, binary)
        binary.chmod(0o755)
        _VISION_HELPER_BINARY = binary
        _VISION_HELPER_READY = True
        _log_processor_event("DEBUG", "Vision OCR helper available", extractor="vision", binary=str(binary))
        return str(binary)
    except Exception as e:
        _VISION_HELPER_READY = False
        _VISION_HELPER_BINARY = None
        _log_processor_event(
            "WARN",
            "Vision OCR helper compile failed",
            extractor="vision",
            error=str(e),
        )
        return None


def _reconstruct_column_aware_text(observations: list[dict]) -> str:
    """
    Group observations by approximate row (Y-band), detect columns by X,
    then output left-to-right within each row. Preserves spatial structure of
    multi-column tax forms (W-2, 1099) better than naive top-to-bottom sort.
    """
    if not observations:
        return ""
    # Sort by Y descending (top of page first — Vision uses bottom-left origin)
    obs = sorted(observations, key=lambda o: -o["y"])
    # Group into rows: observations within 0.012 Y of each other are the same row
    rows: list[list[dict]] = []
    for o in obs:
        if rows and abs(o["y"] - rows[-1][0]["y"]) < 0.012:
            rows[-1].append(o)
        else:
            rows.append([o])
    # Within each row, sort left-to-right by X
    lines: list[str] = []
    for row in rows:
        row.sort(key=lambda o: o["x"])
        lines.append("  ".join(o["text"] for o in row))
    return "\n".join(lines)


def _ocr_with_vision_payload_path(image_path: str) -> dict[str, Any] | None:
    """Run OCR with Apple's Vision framework when available and return raw payload."""
    helper = _ensure_vision_helper_binary()
    if helper is None:
        return None
    try:
        proc = subprocess.run(
            [helper, image_path],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if proc.returncode != 0:
            _log_processor_event(
                "WARN",
                "Vision OCR failed",
                extractor="vision",
                error=(proc.stderr or proc.stdout or "").strip()[:400],
                image_path=image_path,
            )
            return None
        payload = json.loads((proc.stdout or "{}").strip() or "{}")
        if not isinstance(payload, dict):
            return None
        return payload
    except Exception as e:
        _log_processor_event(
            "WARN",
            "Vision OCR invocation failed",
            extractor="vision",
            error=str(e),
            image_path=image_path,
        )
        return None


def _ocr_with_vision_path(image_path: str) -> str | None:
    payload = _ocr_with_vision_payload_path(image_path)
    if not payload:
        return None
    observations = payload.get("observations")
    if isinstance(observations, list) and observations:
        text = _reconstruct_column_aware_text(observations)
    else:
        text = str(payload.get("text") or "").strip()
    return text or None


def _summarize_image_with_vision_model(image_bytes: bytes, source_path: str, visible_text: str) -> tuple[str, str]:
    """Generate a compact ingest-time summary for image retrieval."""
    model = ensure_vision_model_ready()
    try:
        import ollama

        visible_section = ""
        cleaned_visible = (visible_text or "").strip()
        if cleaned_visible:
            preview = cleaned_visible[:1200]
            visible_section = f"\nVisible text extracted from OCR:\n{preview}\n"
        resp = ollama.chat(
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": (
                        "Summarize this image for local retrieval in 1-2 concise sentences. "
                        "Mention the main subject, scene, and any obvious UI/app context. "
                        "Include obvious colors, physical traits, or other salient visual attributes when clearly visible. "
                        "Do not invent unreadable text or identities."
                        f"{visible_section}"
                    ),
                    "images": [image_bytes],
                }
            ],
            keep_alive=0,
            options={"temperature": 0, "seed": 42},
        )
        text = ((resp.get("message") or {}).get("content") or "").strip()
    except Exception as e:
        raise ImageExtractionError(f"Vision image summary failed for {source_path}: {e}") from e
    if not text:
        raise ImageExtractionError(f"Vision image summary returned no content for {source_path}.")
    return text, model


def _merge_observation_rows(observations: list[dict[str, Any]]) -> list[ImageRegion]:
    """Merge Vision OCR observations into small semantic text regions."""
    if not observations:
        return []
    obs = [
        {
            "text": str(o.get("text") or "").strip(),
            "x": float(o.get("x") or 0.0),
            "y": float(o.get("y") or 0.0),
            "w": float(o.get("w") or 0.0),
            "h": float(o.get("h") or 0.0),
        }
        for o in observations
        if str(o.get("text") or "").strip()
    ]
    if not obs:
        return []
    obs.sort(key=lambda o: (-o["y"], o["x"]))
    median_h = sorted(o["h"] for o in obs)[len(obs) // 2] if obs else 0.03
    row_threshold = max(0.012, median_h * 0.7)
    rows: list[list[dict[str, Any]]] = []
    for item in obs:
        if rows and abs(item["y"] - rows[-1][0]["y"]) <= row_threshold:
            rows[-1].append(item)
        else:
            rows.append([item])

    row_entries: list[dict[str, Any]] = []
    for row in rows:
        row.sort(key=lambda o: o["x"])
        x0 = min(o["x"] for o in row)
        y0 = min(o["y"] for o in row)
        x1 = max(o["x"] + o["w"] for o in row)
        y1 = max(o["y"] + o["h"] for o in row)
        row_entries.append(
            {
                "text": " ".join(o["text"] for o in row).strip(),
                "x0": x0,
                "y0": y0,
                "x1": x1,
                "y1": y1,
                "h": y1 - y0,
            }
        )

    blocks: list[list[dict[str, Any]]] = []
    for row in row_entries:
        if not blocks:
            blocks.append([row])
            continue
        prev = blocks[-1][-1]
        vertical_gap = prev["y0"] - row["y1"]
        overlap = min(prev["x1"], row["x1"]) - max(prev["x0"], row["x0"])
        min_width = max(0.0001, min(prev["x1"] - prev["x0"], row["x1"] - row["x0"]))
        overlap_ratio = overlap / min_width if overlap > 0 else 0.0
        left_shift = abs(prev["x0"] - row["x0"])
        if vertical_gap <= max(0.03, max(prev["h"], row["h"]) * 1.6) and (overlap_ratio >= 0.18 or left_shift <= 0.08):
            blocks[-1].append(row)
        else:
            blocks.append([row])

    regions: list[ImageRegion] = []
    for block in blocks:
        text = "\n".join(r["text"] for r in block if r["text"]).strip()
        if not text:
            continue
        x0 = min(r["x0"] for r in block)
        y0 = min(r["y0"] for r in block)
        x1 = max(r["x1"] for r in block)
        y1 = max(r["y1"] for r in block)
        regions.append(
            ImageRegion(
                text=text,
                role="ocr_block",
                x=round(x0, 6),
                y=round(y0, 6),
                w=round(max(0.0, x1 - x0), 6),
                h=round(max(0.0, y1 - y0), 6),
                needs_vision_enrichment=False,
            )
        )
    return regions


def _ocr_with_paddle_path(image_path: str) -> str | None:
    """Run OCR with PaddleOCR when available. Returns extracted text or None."""
    engine = _get_paddle_ocr_engine()
    if engine is None:
        return None
    try:
        use_angle = _ocr_preprocess_enabled()
        result = engine.ocr(str(image_path), cls=use_angle)
        text = _extract_paddleocr_text(result)
        return text or None
    except Exception as e:
        _log_processor_event(
            "WARN",
            "PaddleOCR inference failed",
            error=str(e),
            extractor="paddleocr",
        )
        return None


def _ocr_with_paddle_detail_path(image_path: str) -> _OCRResult | None:
    text = _ocr_with_paddle_path(image_path)
    if not text:
        return None
    return _OCRResult(text=text, backend="paddleocr")


def _ocr_with_paddle(image_png: bytes) -> str | None:
    """Backwards-compatible byte wrapper around PaddleOCR."""
    try:
        with tempfile.TemporaryDirectory() as td:
            image_path = Path(td) / "page.png"
            image_path.write_bytes(image_png)
            return _ocr_with_paddle_path(str(image_path))
    except Exception as e:
        _log_processor_event(
            "WARN",
            "PaddleOCR inference failed",
            error=str(e),
            extractor="paddleocr",
        )
        return None


def _ocr_with_tesseract_path(image_path: str) -> str | None:
    """Run OCR with tesseract CLI when available. Returns extracted text or None."""
    if not _tesseract_available():
        return None
    try:
        def _run_tesseract(psm: str) -> str | None:
            with tempfile.TemporaryDirectory() as td:
                output_base = Path(td) / f"ocr_out_psm{psm}"
                cmd = [
                    "tesseract",
                    str(image_path),
                    str(output_base),
                    "-l",
                    "eng",
                    "--psm",
                    psm,
                ]
                proc = subprocess.run(cmd, capture_output=True, text=True)
                if proc.returncode != 0:
                    _log_processor_event(
                        "WARN",
                        "tesseract OCR failed",
                        error=(proc.stderr or "").strip()[:240],
                        extractor="tesseract",
                        psm=psm,
                    )
                    return None
                out_path = output_base.with_suffix(".txt")
                if not out_path.exists():
                    return None
                text = out_path.read_text(encoding="utf-8", errors="replace").strip()
                return text or None

        if not _ocr_preprocess_enabled():
            return _run_tesseract("6")

        text_6 = _run_tesseract("6")
        text_11 = _run_tesseract("11")
        ratio_6 = _alnum_ratio(text_6 or "")
        ratio_11 = _alnum_ratio(text_11 or "")

        chosen_psm = "6"
        chosen_text = text_6
        if text_11 and (not chosen_text or ratio_11 > ratio_6):
            chosen_psm = "11"
            chosen_text = text_11

        _log_processor_event(
            "DEBUG",
            "tesseract OCR PSM comparison",
            psm_6_ratio=round(ratio_6, 4),
            psm_11_ratio=round(ratio_11, 4),
            chosen_psm=chosen_psm if chosen_text else None,
        )
        return chosen_text
    except Exception as e:
        _log_processor_event(
            "WARN",
            "tesseract OCR invocation failed",
            error=str(e),
            extractor="tesseract",
        )
        return None


def _ocr_with_tesseract_detail_path(image_path: str) -> _OCRResult | None:
    text = _ocr_with_tesseract_path(image_path)
    if not text:
        return None
    return _OCRResult(text=text, backend="tesseract")


def _ocr_with_tesseract(image_png: bytes) -> str | None:
    """Backwards-compatible byte wrapper around tesseract."""
    try:
        with tempfile.TemporaryDirectory() as td:
            image_path = Path(td) / "page.png"
            image_path.write_bytes(image_png)
            return _ocr_with_tesseract_path(str(image_path))
    except Exception as e:
        _log_processor_event(
            "WARN",
            "tesseract OCR invocation failed",
            error=str(e),
            extractor="tesseract",
        )
        return None


def _ocr_with_vision_detail_path(image_path: str) -> _OCRResult | None:
    payload = _ocr_with_vision_payload_path(image_path)
    if not payload:
        return None
    observations = payload.get("observations") if isinstance(payload.get("observations"), list) else []
    text = _reconstruct_column_aware_text(observations) if observations else str(payload.get("text") or "").strip()
    if not text:
        return None
    return _OCRResult(
        text=text,
        backend="vision",
        observations=tuple(observations),
        raw_payload=payload,
    )


def _ocr_image_path_detailed(image_path: str, source_path: str, ocr_mode: str) -> _OCRResult | None:
    configured = _configured_ocr_backend()
    preferred = _preferred_ocr_backends()
    for idx, backend in enumerate(preferred):
        if not _ocr_backend_available(backend):
            level = "WARN" if configured != "auto" else "DEBUG"
            _log_processor_event(
                level,
                "OCR backend unavailable",
                path=source_path,
                backend=backend,
                ocr_mode=ocr_mode,
            )
            if configured != "auto":
                return None
            continue

        candidate: _OCRResult | None = None
        if backend == "vision":
            candidate = _ocr_with_vision_detail_path(image_path)
        elif backend == "paddleocr":
            candidate = _ocr_with_paddle_detail_path(image_path)
        elif backend == "tesseract":
            candidate = _ocr_with_tesseract_detail_path(image_path)

        if candidate and candidate.text:
            quality_ok, reasons, quality_stats = _ocr_quality_assessment(candidate.text)
            if not quality_ok:
                _log_processor_event(
                    "INFO",
                    "OCR text dropped by quality gate",
                    path=source_path,
                    backend=backend,
                    ocr_mode=ocr_mode,
                    reasons=list(reasons),
                    **quality_stats,
                )
                candidate = None
            else:
                candidate = _OCRResult(
                    text=candidate.text,
                    backend=candidate.backend,
                    observations=candidate.observations,
                    raw_payload=candidate.raw_payload,
                    quality_reasons=reasons,
                )
        if candidate and candidate.text:
            _log_processor_event(
                "INFO",
                "OCR backend selected",
                path=source_path,
                backend=backend,
                ocr_mode=ocr_mode,
            )
            return candidate

        level = "WARN" if configured != "auto" else "INFO"
        _log_processor_event(
            level,
            "OCR backend produced no text",
            path=source_path,
            backend=backend,
            ocr_mode=ocr_mode,
        )
        if configured == "auto":
            for next_backend in preferred[idx + 1 :]:
                if _ocr_backend_available(next_backend):
                    _log_processor_event(
                        "INFO",
                        "OCR backend fallback",
                        path=source_path,
                        from_backend=backend,
                        to_backend=next_backend,
                        ocr_mode=ocr_mode,
                    )
                    break
            continue
        return None
    return None


def _ocr_image_path(image_path: str, source_path: str, ocr_mode: str) -> tuple[str | None, str | None]:
    result = _ocr_image_path_detailed(image_path, source_path, ocr_mode)
    if result is None:
        return None, None
    return result.text, result.backend


def _ocr_image_bytes(
    image_bytes: bytes,
    source_path: str,
    ocr_mode: str,
    preferred_suffix: str = ".png",
) -> tuple[str | None, str | None]:
    result = _ocr_image_bytes_detailed(
        image_bytes,
        source_path,
        ocr_mode=ocr_mode,
        preferred_suffix=preferred_suffix,
    )
    if result is None:
        return None, None
    return result.text, result.backend


def _ocr_image_bytes_detailed(
    image_bytes: bytes,
    source_path: str,
    ocr_mode: str,
    preferred_suffix: str = ".png",
) -> _OCRResult | None:
    payload = _preprocess_image_for_ocr(image_bytes) if _ocr_preprocess_enabled() else image_bytes
    suffix = preferred_suffix if preferred_suffix.startswith(".") else f".{preferred_suffix}"
    try:
        with tempfile.TemporaryDirectory() as td:
            image_path = Path(td) / f"ocr-input{suffix}"
            image_path.write_bytes(payload)
            return _ocr_image_path_detailed(str(image_path), source_path, ocr_mode)
    except Exception as e:
        _log_processor_event(
            "WARN",
            "OCR image staging failed",
            path=source_path,
            error=str(e),
            ocr_mode=ocr_mode,
        )
        return None


def _ocr_image_file(data: bytes, source_path: str, ocr_mode: str = "image_file") -> tuple[str | None, str | None]:
    result = _ocr_image_file_detailed(data, source_path, ocr_mode=ocr_mode)
    if result is None:
        return None, None
    return result.text, result.backend


def _ocr_image_file_detailed(data: bytes, source_path: str, ocr_mode: str = "image_file") -> _OCRResult | None:
    suffix = Path(source_path).suffix.lower() or ".png"
    if _ocr_preprocess_enabled():
        return _ocr_image_bytes_detailed(data, source_path, ocr_mode=ocr_mode, preferred_suffix=".png")
    if Path(source_path).exists():
        return _ocr_image_path_detailed(source_path, source_path, ocr_mode)
    return _ocr_image_bytes_detailed(data, source_path, ocr_mode=ocr_mode, preferred_suffix=suffix)


def _ocr_pdf_page_text(page: Any, source_path: str) -> tuple[str | None, str | None]:
    """OCR fallback for image-only PDF pages. Returns (text, backend) or (None, None)."""
    result = _ocr_pdf_page_result(page, source_path)
    if result is None:
        return None, None
    return result.text, result.backend


def _ocr_pdf_page_result(page: Any, source_path: str) -> _OCRResult | None:
    """OCR fallback for image-only PDF pages with quality gate metadata."""
    use_preprocess = _ocr_preprocess_enabled()
    try:
        dpi = 400 if use_preprocess else 300
        image_png = page.get_pixmap(dpi=dpi, alpha=False).tobytes("png")
    except Exception as e:
        _log_processor_event(
            "WARN",
            "PDF OCR render failed",
            path=source_path,
            page=(getattr(page, "number", 0) + 1),
            error=str(e),
        )
        return None

    return _ocr_image_bytes_detailed(image_png, source_path, ocr_mode="pdf_scan_fallback", preferred_suffix=".png")


def _normalize_table_rows(rows: list[list[str | None]]) -> list[list[str]]:
    """Normalize table rows: strip/compact whitespace, trim trailing empty cells, and drop empty rows."""
    out: list[list[str]] = []
    for row in rows or []:
        normalized = [re.sub(r"\s+", " ", (cell or "").strip()) for cell in (row or [])]
        first_non_empty = -1
        last_non_empty = -1
        for i, cell in enumerate(normalized):
            if cell:
                if first_non_empty < 0:
                    first_non_empty = i
                last_non_empty = i
        if first_non_empty >= 0 and last_non_empty >= 0:
            out.append(normalized[first_non_empty : last_non_empty + 1])
    return out


def _table_to_markdown(rows: list[list[str]]) -> str:
    """Convert normalized rows to markdown table text."""
    if not rows:
        return ""
    col_count = max(len(r) for r in rows)

    def _pad(row: list[str]) -> list[str]:
        cells = row + ([""] * (col_count - len(row)))
        return [c.replace("|", "\\|") for c in cells]

    header = _pad(rows[0])
    body = [_pad(r) for r in rows[1:]]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * col_count) + " |",
    ]
    lines.extend("| " + " | ".join(r) + " |" for r in body)
    return "\n".join(lines)


def _extract_line_value_hints(rows: list[list[str]]) -> list[str]:
    """Extract compact line/value hints from table-like rows (e.g., line 9 -> 7,522.)."""
    def _is_high_signal_numeric(value: str) -> bool:
        if not _NUMERIC_VALUE_PATTERN.match(value):
            return False
        digits_only = re.sub(r"\D", "", value)
        # Ignore tiny integers that are usually row/column identifiers.
        if len(digits_only) <= 2 and "$" not in value and "," not in value and "." not in value:
            return False
        return True

    hints: list[str] = []
    seen: set[str] = set()
    for row in rows:
        if not row:
            continue
        for idx, cell in enumerate(row):
            token = cell.strip().lower()
            if not token:
                continue
            m = _LINE_TOKEN_PATTERN.match(token)
            if not m:
                continue
            line_label = m.group(1)
            num_match = re.match(r"(\d+)", line_label)
            if num_match and int(num_match.group(1)) > _MAX_HINT_LINE_NUMBER:
                break
            value = ""
            for j in range(idx + 1, len(row)):
                candidate = row[j].strip()
                if candidate and _is_high_signal_numeric(candidate):
                    value = candidate
                    break
            if value:
                hint = f"line {line_label}: {value}"
                if hint not in seen:
                    seen.add(hint)
                    hints.append(hint)
            break
    return hints


def _merge_pdf_page_content(
    raw_text: str,
    table_md: str,
    hints: list[str],
    ocr_text: str | None = None,
    max_structured_chars: int = _PDF_STRUCTURED_MAX_CHARS,
) -> str:
    """Merge structured table context ahead of raw text with a hard cap on structured size."""
    sections: list[str] = []
    if hints:
        sections.append("Structured values:\n" + "\n".join(hints))
    if table_md:
        sections.append("Extracted tables:\n" + table_md)
    if ocr_text:
        sections.append(f"{_OCR_LABEL}\n{ocr_text.strip()}")
    structured = "\n\n".join(s for s in sections if s).strip()
    if structured and len(structured) > max_structured_chars:
        structured = structured[:max_structured_chars].rstrip()
    raw = (raw_text or "").strip()
    if structured and raw:
        return f"{structured}\n\n{raw}"
    return structured or raw


def _extract_pdf_tables_by_page(data: bytes) -> list[list[list[str | None]]]:
    """Return page-indexed tables from pdfplumber. Empty list means unavailable/no tables."""
    try:
        import pdfplumber
    except ImportError:
        return []
    # pdfplumber/pdfminer can emit high-volume parser warnings on broken PDFs; keep stderr usable.
    logging.getLogger("pdfminer").setLevel(logging.ERROR)

    tables_by_page: list[list[list[str | None]]] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            page_tables = page.extract_tables() or []
            tables_by_page.append(page_tables)
    return tables_by_page


class PDFProcessor:
    format_label = "PDF"
    install_hint = "pymupdf"

    def extract(self, data: bytes, source_path: str) -> list[ExtractedPage]:
        try:
            import fitz
            tables_by_page: list[list[list[str | None]]] = []
            if _pdf_tables_enabled():
                try:
                    tables_by_page = _extract_pdf_tables_by_page(data)
                except Exception as e:
                    _log_processor_event(
                        "WARN",
                        "PDF table extraction failed",
                        path=source_path,
                        error=str(e),
                        extractor="pdfplumber",
                    )
            with fitz.open(stream=data, filetype="pdf") as doc:
                out: list[ExtractedPage] = []
                pages_without_text_or_ocr: list[int] = []
                for page in doc:
                    raw_text = page.get_text()
                    ocr_text: str | None = None
                    page_meta: dict[str, Any] = {}
                    try:
                        for w in page.widgets():
                            name = getattr(w, "field_name", None) or ""
                            val = getattr(w, "field_value", None)
                            if name and val is not None and str(val).strip():
                                raw_text += f"\n{name}: {val}"
                    except Exception:
                        pass
                    if not raw_text.strip():
                        ocr_text, backend = _ocr_pdf_page_text(page, source_path)
                        if ocr_text:
                            if backend:
                                page_meta = {"ocr_backend": backend, "ocr_mode": "pdf_scan_fallback"}
                            _log_processor_event(
                                "INFO",
                                "PDF OCR fallback applied",
                                path=source_path,
                                page=page.number + 1,
                                backend=backend,
                            )
                        else:
                            pages_without_text_or_ocr.append(page.number + 1)
                    page_tables = tables_by_page[page.number] if page.number < len(tables_by_page) else []
                    hints: list[str] = []
                    md_tables: list[str] = []
                    for table in page_tables:
                        normalized = _normalize_table_rows(table)
                        if not normalized:
                            continue
                        hints.extend(_extract_line_value_hints(normalized))
                        md = _table_to_markdown(normalized)
                        if md:
                            md_tables.append(md)
                    text = _merge_pdf_page_content(raw_text, "\n\n".join(md_tables), hints, ocr_text=ocr_text)
                    out.append(ExtractedPage(text=text, page_num=page.number + 1, meta=page_meta or None))
                if pages_without_text_or_ocr:
                    _log_processor_event(
                        "WARN",
                        "PDF pages have no extractable text and OCR produced no text",
                        path=source_path,
                        pages=pages_without_text_or_ocr,
                        page_count=len(pages_without_text_or_ocr),
                        available_ocr_backends=_available_ocr_backends(),
                    )
                return out
        except Exception as e:
            raise PDFExtractionError(str(e)) from e


class DOCXProcessor:
    format_label = "Document"
    install_hint = "python-docx"

    def extract(self, data: bytes, source_path: str) -> str | None:
        try:
            import docx
            doc = docx.Document(io.BytesIO(data))
            return "\n".join(para.text for para in doc.paragraphs)
        except ImportError:
            return None
        except Exception as e:
            raise DOCXExtractionError(str(e)) from e


class XLSXProcessor:
    format_label = "Spreadsheet"
    install_hint = "openpyxl"

    def extract(self, data: bytes, source_path: str) -> str | None:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join(str(c) if c is not None else "" for c in row)
                    if line.strip():
                        parts.append(line)
            wb.close()
            text = "\n".join(parts)
            return text.strip() or None
        except ImportError:
            return None
        except Exception:
            raise XLSXExtractionError("Failed to extract XLSX content.")


class PPTXProcessor:
    format_label = "Presentation"
    install_hint = "python-pptx"

    def extract(self, data: bytes, source_path: str) -> str | None:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"Slide {i + 1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            text = "\n".join(parts)
            return text.strip() or None
        except ImportError:
            return None
        except Exception:
            raise PPTXExtractionError("Failed to extract PPTX content.")


def _build_image_summary_text(summary: str, visible_text: str) -> str:
    lines = [f"Image summary: {(summary or '').strip()}"]
    cleaned_visible = (visible_text or "").strip()
    if cleaned_visible:
        preview_lines = [line.strip() for line in cleaned_visible.splitlines() if line.strip()]
        preview = "\n".join(preview_lines[:12]).strip()
        if preview:
            lines.append("Visible text:")
            lines.append(preview[:1600])
    return "\n".join(lines).strip()


class ImageProcessor:
    format_label = "Image"
    install_hint = "macOS Vision or paddleocr/tesseract with LLMLIBRARIAN_VISION_MODEL"

    def extract(
        self,
        data: bytes,
        source_path: str,
        *,
        enable_multimodal: bool = True,
    ) -> ExtractedImage | None:
        try:
            ocr_result = _ocr_image_file_detailed(data, source_path, ocr_mode="image_file")
            signal = _image_ocr_signal_assessment(ocr_result)
            raw_visible_text = (ocr_result.text if ocr_result else "").strip()
            visible_text = raw_visible_text if signal["keep_visible_text"] else ""
            backend = ocr_result.backend if ocr_result else None

            regions: list[ImageRegion] = []
            raw_payload: dict[str, Any] | None = None
            if ocr_result and ocr_result.backend == "vision" and ocr_result.observations and signal["keep_visible_text"]:
                regions = _merge_observation_rows(list(ocr_result.observations))
                raw_payload = dict(ocr_result.raw_payload or {})
            elif visible_text:
                regions = [
                    ImageRegion(
                        text=visible_text,
                        role="ocr_block",
                        x=0.0,
                        y=0.0,
                        w=1.0,
                        h=1.0,
                        needs_vision_enrichment=False,
                    )
                ]

            summary_text = ""
            vision_model = ""
            summary_status = "deferred"
            if enable_multimodal:
                if signal["eager_summary"]:
                    summary_text, vision_model = _summarize_image_with_vision_model(data, source_path, visible_text)
                    summary_status = "eager"
                else:
                    summary_text = _image_summary_placeholder(visible_text)
            else:
                summary_text = _image_summary_disabled_placeholder(visible_text)
                summary_status = "disabled"

            if not summary_text:
                summary_text = visible_text or Path(source_path).name

            if not regions:
                regions = [
                    ImageRegion(
                        text=summary_text,
                        role="full_frame_summary",
                        x=0.0,
                        y=0.0,
                        w=1.0,
                        h=1.0,
                        needs_vision_enrichment=(summary_status == "deferred"),
                    )
                ]

            artifact = {
                "source_path": source_path,
                "ocr_backend": backend,
                "ocr_mode": "image_file",
                "visible_text": visible_text,
                "raw_visible_text": raw_visible_text,
                "summary": summary_text,
                "summary_status": summary_status,
                "ocr_signal_score": signal["ocr_signal_score"],
                "query_time_summary_cached_at": None,
                "ocr_signal": {
                    "observation_count": signal["observation_count"],
                    "coverage": signal["coverage"],
                    "quality_ok": signal["quality_ok"],
                    "quality_reasons": list(signal["quality_reasons"]),
                    **dict(signal["quality_stats"]),
                },
                "vision_model": vision_model or None,
                "regions": [
                    {
                        "text": region.text,
                        "role": region.role,
                        "x": region.x,
                        "y": region.y,
                        "w": region.w,
                        "h": region.h,
                        "needs_vision_enrichment": region.needs_vision_enrichment,
                        "meta": dict(region.meta or {}),
                    }
                    for region in regions
                ],
                "raw_vision_payload": raw_payload,
            }
            meta = {
                "ocr_backend": backend,
                "ocr_mode": "image_file",
                "source_modality": "image",
                "summary_status": summary_status,
                "needs_vision_enrichment": summary_status == "deferred",
            }
            if vision_model:
                meta["vision_model"] = vision_model
            return ExtractedImage(
                summary=_build_image_summary_text(summary_text, visible_text),
                visible_text=visible_text,
                regions=tuple(regions),
                meta=meta,
                artifact=artifact,
            )
        except Exception as e:
            raise ImageExtractionError(str(e)) from e


class TextProcessor:
    format_label = "Text"
    install_hint = ""

    def extract(self, data: bytes, source_path: str) -> str:
        try:
            return data.decode("utf-8", errors="replace")
        except Exception as e:
            raise TextExtractionError(str(e)) from e


# Registry: suffix -> processor instance
PROCESSORS: dict[str, DocumentProcessor] = {
    ".pdf": PDFProcessor(),  # type: ignore[dict-item]
    ".docx": DOCXProcessor(),  # type: ignore[dict-item]
    ".xlsx": XLSXProcessor(),  # type: ignore[dict-item]
    ".pptx": PPTXProcessor(),  # type: ignore[dict-item]
    ".png": ImageProcessor(),  # type: ignore[dict-item]
    ".jpg": ImageProcessor(),  # type: ignore[dict-item]
    ".jpeg": ImageProcessor(),  # type: ignore[dict-item]
    ".heic": ImageProcessor(),  # type: ignore[dict-item]
    ".heif": ImageProcessor(),  # type: ignore[dict-item]
    ".tif": ImageProcessor(),  # type: ignore[dict-item]
    ".tiff": ImageProcessor(),  # type: ignore[dict-item]
}

# Default processor for code/text files
DEFAULT_PROCESSOR = TextProcessor()
